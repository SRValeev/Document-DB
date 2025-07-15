import os
import re
import json
import fitz
import spacy
import pytesseract
import logging
import camelot
import pandas as pd
from pathlib import Path
from PIL import Image
from io import BytesIO
from datetime import datetime
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from tqdm import tqdm
from typing import List, Dict, Optional, Union, Callable
from sentence_transformers import SentenceTransformer
from utils.helpers import load_config, normalize_text, generate_unique_id
import warnings

warnings.filterwarnings('ignore', category=pd.errors.SettingWithCopyWarning)
warnings.filterwarnings('ignore', category=FutureWarning)

class FileProcessor:
    def __init__(self, config):
        self.config = config
        self.logger = logging.getLogger(__name__)
        self._init_models()
        self._init_processing_params()
        self._init_regex_patterns()
        
    def _init_models(self):
        """Инициализация моделей обработки"""
        try:
            self.nlp = spacy.load(
                self.config['processing']['spacy_model'],
                disable=["parser", "lemmatizer", "ner"]
            )
            self.nlp.add_pipe('sentencizer')
            
            self.embedding_model = SentenceTransformer(
                self.config['processing']['embedding_model'],
                device=self.config['processing']['device']
            )
            self.logger.info("Models loaded successfully")
        except Exception as e:
            self.logger.error(f"Model loading error: {str(e)}", exc_info=True)
            raise

    def _init_processing_params(self):
        """Инициализация параметров обработки из конфига"""
        processing_cfg = self.config['processing']
        self.chunk_size = processing_cfg['chunk_size']
        self.chunk_overlap = processing_cfg['chunk_overlap']
        self.min_chunk_size = processing_cfg['min_chunk_size']
        self.smart_chunking = processing_cfg['smart_chunking']
        self.min_similarity = processing_cfg['min_similarity']
        self.use_ocr = self.config['ocr']['enabled']
        self.ocr_languages = "+".join(self.config['ocr']['languages'])
        self.extract_tables = self.config['tables']['enabled']
        self.max_table_size = self.config['tables']['max_table_size']
        self.table_format = self.config['tables']['format']
        self.max_file_size_mb = processing_cfg['max_file_size_mb']

    def _init_regex_patterns(self):
        """Инициализация regex-паттернов из конфига"""
        patterns = self.config['processing'].get('regex_patterns', {})
        self.header_pattern = re.compile(
            patterns.get('header', r'^(Глава|Раздел|Часть|Параграф)\s+\d+[.:]?\s+'),
            re.IGNORECASE
        )
        self.subheader_pattern = re.compile(
            patterns.get('subheader', r'^\d+\.\d+\.\s+'),
            re.IGNORECASE
        )

    def process_file(self, file_path: str) -> List[Dict]:
        """Основной метод обработки файла с проверкой размера"""
        if not os.path.exists(file_path):
            self.logger.error(f"File not found: {file_path}")
            return []

        file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
        if file_size_mb > self.max_file_size_mb:
            self.logger.warning(f"File {file_path} exceeds max size {self.max_file_size_mb}MB")
            return []

        file_ext = os.path.splitext(file_path)[1].lower()
        processor = self._get_processor(file_ext)
        
        if not processor:
            self.logger.warning(f"Unsupported format: {file_path}")
            return []

        file_id = self._generate_file_id(file_path)
        return processor(file_path, file_id)

    def _get_processor(self, file_ext: str) -> Optional[Callable]:
        """Возвращает обработчик для конкретного формата файла"""
        processors = {
            '.pdf': self._process_pdf,
            '.doc': self._process_docx,
            '.docx': self._process_docx,
            '.pptx': self._process_pptx,
            '.xlsx': self._process_xlsx,
            '.txt': self._process_text
        }
        return processors.get(file_ext)

    def _generate_file_id(self, file_path: str) -> str:
        """Генерация уникального ID для файла"""
        file_name = os.path.basename(file_path)
        if '_' in file_name:
            original_name = '_'.join(file_name.split('_')[:-1]) + Path(file_name).suffix
        else:
            original_name = file_name
        return f"{original_name}_{generate_unique_id()}"

    def _process_pdf(self, file_path: str, file_id: str) -> List[Dict]:
        """Обработка PDF файлов"""
        chunks = []
        current_chapter = ""
        current_section = ""

        with fitz.open(file_path) as doc:
            toc = doc.get_toc() if hasattr(doc, 'get_toc') else []
            
            for page_num in tqdm(range(len(doc)), desc=f"Processing PDF {os.path.basename(file_path)}"):
                page = doc.load_page(page_num)
                current_chapter, current_section = self._update_sections_from_toc(toc, page_num, current_chapter, current_section)
                
                # Обработка текста
                text = page.get_text("text")
                if text.strip():
                    chunks.extend(self._process_text_content(
                        text, file_id, page_num, "text", current_chapter, current_section
                    ))

                # Обработка изображений (OCR)
                if self.use_ocr:
                    chunks.extend(self._process_pdf_images(page, file_id, page_num))

                # Извлечение таблиц
                if self.extract_tables:
                    chunks.extend(self._extract_tables_from_pdf(file_path, page_num+1, file_id))

        return chunks

    def _update_sections_from_toc(self, toc: List, page_num: int, current_chapter: str, current_section: str) -> tuple:
        """Обновляет текущие разделы на основе оглавления"""
        for item in toc:
            if item[2] == page_num + 1:
                if item[0] == 1:
                    current_chapter = normalize_text(item[1])
                elif item[0] == 2:
                    current_section = normalize_text(item[1])
        return current_chapter, current_section

    def _process_docx(self, file_path: str, file_id: str) -> List[Dict]:
        """Обработка DOCX файлов"""
        chunks = []
        current_chapter = ""
        current_section = ""
        
        try:
            doc = Document(file_path)
            full_text = []

            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                    
                style = para.style.name.lower() if para.style else ""
                if 'heading 1' in style:
                    current_chapter = text
                elif 'heading 2' in style:
                    current_section = text
                
                full_text.append(text)

            if full_text:
                chunks.extend(self._process_text_content(
                    "\n".join(full_text), file_id, 0, "text", current_chapter, current_section
                ))

            # Обработка таблиц
            if self.extract_tables:
                for table in doc.tables:
                    chunks.extend(self._extract_tables_from_docx(table, file_id, current_chapter, current_section))

        except Exception as e:
            self.logger.error(f"Error processing DOCX {file_path}: {str(e)}")

        return chunks

    def _process_pptx(self, file_path: str, file_id: str) -> List[Dict]:
        """Обработка PPTX файлов"""
        chunks = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                
                if slide_text:
                    chunks.extend(self._process_text_content(
                        "\n".join(slide_text), file_id, slide_num, "slide", "", ""
                    ))
                    
        except Exception as e:
            self.logger.error(f"Error processing PPTX {file_path}: {str(e)}")
            
        return chunks

    def _process_xlsx(self, file_path: str, file_id: str) -> List[Dict]:
        """Обработка XLSX файлов"""
        chunks = []
        
        try:
            wb = load_workbook(file_path, read_only=True)
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                table_data = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    table_data.append(row_data)
                
                if table_data:
                    df = pd.DataFrame(table_data[1:], columns=table_data[0])
                    table_text = self._format_table(df, sheet_name, 0)
                    chunks.append(self._create_chunk(
                        text=table_text,
                        file_id=file_id,
                        page=0,
                        content_type="table",
                        chapter=f"Sheet: {sheet_name}",
                        section="",
                        chunk_order=0
                    ))
                    
        except Exception as e:
            self.logger.error(f"Error processing XLSX {file_path}: {str(e)}")
            
        return chunks

    def _process_text(self, file_path: str, file_id: str) -> List[Dict]:
        """Обработка текстовых файлов"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            if text.strip():
                return self._process_text_content(text, file_id, 0, "text", "", "")
        except Exception as e:
            self.logger.error(f"Error processing text file {file_path}: {str(e)}")
        
        return []

    def _process_text_content(self, text: str, file_id: str, page: int, 
                            content_type: str, chapter: str, section: str) -> List[Dict]:
        """Обработка текстового контента с разделением на чанки"""
        if self.smart_chunking:
            return self._split_text_into_chunks(text, file_id, page, content_type, chapter, section)
        else:
            return [self._create_chunk(
                text=text,
                file_id=file_id,
                page=page,
                content_type=content_type,
                chapter=chapter,
                section=section,
                chunk_order=0
            )]

    def _split_text_into_chunks(self, text: str, file_id: str, page: int, 
                              content_type: str, chapter: str, section: str) -> List[Dict]:
        """Умное разделение текста на чанки с контекстом"""
        doc = self.nlp(text)
        chunks = []
        current_chunk = []
        current_length = 0
        
        for sent in doc.sents:
            sent_text = sent.text.strip()
            if not sent_text:
                continue
                
            sent_length = len(sent_text.split())
            
            if current_length + sent_length > self.chunk_size and current_chunk:
                self._save_chunk(current_chunk, file_id, page, content_type, chapter, section, chunks)
                overlap_size = min(self.chunk_overlap, len(current_chunk))
                current_chunk = current_chunk[-overlap_size:]
                current_length = sum(len(s.split()) for s in current_chunk)
            
            current_chunk.append(sent_text)
            current_length += sent_length
        
        if current_chunk:
            self._save_chunk(current_chunk, file_id, page, content_type, chapter, section, chunks)
        
        return self._merge_small_chunks(chunks, self.min_chunk_size)

    def _save_chunk(self, chunk_texts: List[str], file_id: str, page: int, 
                   content_type: str, chapter: str, section: str, chunks: List[Dict]):
        """Сохраняет чанк с метаданными"""
        chunk_text = ' '.join(chunk_texts)
        chunks.append(self._create_chunk(
            text=chunk_text,
            file_id=file_id,
            page=page,
            content_type=content_type,
            chapter=chapter,
            section=section,
            chunk_order=len(chunks))
        )

    def _merge_small_chunks(self, chunks: List[Dict], min_size: int) -> List[Dict]:
        """Объединяет слишком короткие чанки"""
        if not chunks:
            return []
            
        merged = []
        buffer = [chunks[0]]
        buffer_length = len(chunks[0]['text'].split())
        
        for chunk in chunks[1:]:
            chunk_length = len(chunk['text'].split())
            
            if buffer_length + chunk_length <= min_size:
                buffer.append(chunk)
                buffer_length += chunk_length
            else:
                merged.append(self._merge_chunk_buffer(buffer))
                buffer = [chunk]
                buffer_length = chunk_length
        
        if buffer:
            merged.append(self._merge_chunk_buffer(buffer))
        
        return merged

    def _merge_chunk_buffer(self, chunks: List[Dict]) -> Dict:
        """Объединяет несколько чанков в один"""
        merged_text = ' '.join(chunk['text'] for chunk in chunks)
        first_chunk = chunks[0]
        
        return self._create_chunk(
            text=merged_text,
            file_id=first_chunk['metadata']['file_id'],
            page=first_chunk['metadata']['page'],
            content_type=first_chunk['metadata']['type'],
            chapter=first_chunk['metadata'].get('chapter', ''),
            section=first_chunk['metadata'].get('section', ''),
            chunk_order=first_chunk['metadata']['chunk_order']
        )

    def _extract_tables_from_pdf(self, file_path: str, page: int, file_id: str) -> List[Dict]:
        """Извлечение таблиц из PDF"""
        if not self.extract_tables:
            return []

        chunks = []
        try:
            tables = camelot.read_pdf(
                file_path,
                pages=str(page),
                flavor='stream',
                strip_text='\n',
                suppress_stdout=True
            )
            
            for i, table in enumerate(tables):
                try:
                    df = table.df.copy()
                    if df.empty:
                        continue
                        
                    df = df.map(lambda x: str(x).strip() if pd.notna(x) else "")
                    table_text = self._format_table(df, i+1, page)
                    
                    if table_text:
                        chunks.append(self._create_chunk(
                            text=table_text,
                            file_id=file_id,
                            page=page-1,
                            content_type="table",
                            chapter=f"Table {i+1}",
                            section=f"Page {page}",
                            chunk_order=0
                        ))
                except Exception as e:
                    self.logger.error(f"Error processing table {i+1}: {str(e)}")
        
        except Exception as e:
            self.logger.error(f"Table extraction error: {str(e)}")
        
        return chunks

    def _extract_tables_from_docx(self, table, file_id: str, chapter: str, section: str) -> List[Dict]:
        """Извлечение таблиц из DOCX"""
        try:
            data = []
            for row in table.rows:
                data.append([cell.text.strip() for cell in row.cells])
            
            df = pd.DataFrame(data[1:], columns=data[0])
            table_text = self._format_table(df, 0, 0)
            
            return [self._create_chunk(
                text=table_text,
                file_id=file_id,
                page=0,
                content_type="table",
                chapter=chapter,
                section=section,
                chunk_order=0
            )]
        except Exception as e:
            self.logger.error(f"DOCX table processing error: {str(e)}")
            return []

    def _format_table(self, df: pd.DataFrame, table_num: int, page_num: int) -> str:
        """Форматирование таблицы в указанный формат"""
        if self.table_format == "markdown":
            return self._format_table_markdown(df, table_num, page_num)
        elif self.table_format == "csv":
            return df.to_csv(index=False)
        else:
            return df.to_json(indent=2)

    def _format_table_markdown(self, df: pd.DataFrame, table_num: int, page_num: int) -> str:
        """Форматирование таблицы в Markdown"""
        try:
            if len(df) > self.max_table_size:
                df = df.head(self.max_table_size).copy()
                ellipsis_row = pd.Series(["..."] * len(df.columns), index=df.columns)
                df = pd.concat([df, ellipsis_row.to_frame().T])
            
            headers = [str(col).strip() for col in df.columns]
            rows = []
            
            for _, row in df.iterrows():
                rows.append([str(val).strip() if pd.notna(val) else "" for val in row.values])
            
            header_line = "| " + " | ".join(headers) + " |"
            separator = "| " + " | ".join(["---"] * len(headers)) + " |"
            body_lines = ["| " + " | ".join(row) + " |" for row in rows]
            
            return (
                f"Table {table_num} (page {page_num}):\n" +
                header_line + "\n" +
                separator + "\n" +
                "\n".join(body_lines)
            )
        except Exception as e:
            self.logger.error(f"Table formatting error: {str(e)}")
            return f"Failed to process table {table_num}. Text:\n{df.to_string()}"

    def _process_pdf_images(self, page, file_id: str, page_num: int) -> List[Dict]:
        """Обработка изображений в PDF через OCR"""
        if not self.use_ocr:
            return []

        chunks = []
        image_list = page.get_images(full=True)
        
        for img_index, img in enumerate(image_list, 1):
            try:
                base_image = page.extract_image(img[0])
                image = Image.open(BytesIO(base_image["image"]))
                ocr_text = pytesseract.image_to_string(image, lang=self.ocr_languages)
                
                if ocr_text.strip():
                    chunks.append(self._create_chunk(
                        text=f"Image {img_index}:\n{ocr_text}",
                        file_id=file_id,
                        page=page_num,
                        content_type="image",
                        chunk_order=0
                    ))
            except Exception as e:
                self.logger.warning(f"OCR error: {str(e)}")
        
        return chunks

    def _create_chunk(self, text: str, file_id: str, page: int, content_type: str,
                    chapter: str = "", section: str = "", chunk_order: int = 0) -> Dict:
        """Создает структурированный чанк"""
        try:
            if not isinstance(text, str):
                text = str(text)
        except Exception as e:
            text = f"Error converting text: {str(e)}"
        
        chunk_id = generate_unique_id()
        return {
            "id": chunk_id,
            "text": text,
            "embedding": None,
            "metadata": {
                "file_id": file_id,
                "source": file_id if file_id.endswith(tuple(self.config['processing']['supported_formats'])) else os.path.basename(file_id),
                "page": page,
                "type": content_type,
                "chapter": chapter,
                "section": section,
                "chunk_order": chunk_order,
                "processing_date": datetime.now().strftime('%Y-%m-%d'),
                "text_length": len(text),
                "language": self.config['processing'].get('default_language', 'ru')
            }
        }

    def vectorize_chunks(self, chunks: List[Dict]) -> List[Dict]:
        """Векторизация чанков"""
        if not chunks:
            return []
        
        try:
            texts = [chunk['text'] for chunk in chunks]
            embeddings = self.embedding_model.encode(
                texts,
                batch_size=self.config['performance']['embedding_batch_size'],
                show_progress_bar=False
            )
            
            for chunk, embedding in zip(chunks, embeddings):
                # Проверяем на NaN перед сохранением
                if not any(np.isnan(x) for x in embedding):
                    chunk['embedding'] = embedding.tolist()
                else:
                    self.logger.warning(f"NaN values in embedding for chunk: {chunk['id']}")
                    chunk['embedding'] = None
            
            return chunks
        except Exception as e:
            self.logger.error(f"Vectorization error: {str(e)}", exc_info=True)
            return []

    def save_chunks(self, chunks: List[Dict], output_file: str) -> None:
        """Сохранение чанков в файл"""
        try:
            chunks = self.vectorize_chunks(chunks)
            os.makedirs(os.path.dirname(output_file), exist_ok=True)
            
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(chunks, f, ensure_ascii=False, indent=2)
        except Exception as e:
            self.logger.error(f"Error saving chunks: {str(e)}", exc_info=True)

    def create_global_index(self, index_dir: str) -> None:
        """Создает/обновляет глобальный индекс"""
        try:
            index_data = {
                "files": [],
                "total_chunks": 0,
                "timestamp": datetime.now().isoformat()
            }

            processed_dir = self.config['paths']['output_dir']
            for file in os.listdir(processed_dir):
                if file.endswith('.json') and file != self.config['paths']['global_index_file']:
                    file_path = os.path.join(processed_dir, file)
                    with open(file_path, 'r', encoding='utf-8') as f:
                        chunks = json.load(f)
                    
                    if chunks:
                        first_chunk = chunks[0]
                        index_data['files'].append({
                            "id": first_chunk['metadata']['file_id'],
                            "path": first_chunk['metadata']['source'],
                            "chunks_count": len(chunks),
                            "timestamp": first_chunk['metadata']['processing_date']
                        })
                        index_data['total_chunks'] += len(chunks)

            os.makedirs(index_dir, exist_ok=True)
            index_file = os.path.join(index_dir, self.config['paths']['global_index_file'])
            with open(index_file, 'w', encoding='utf-8') as f:
                json.dump(index_data, f, ensure_ascii=False, indent=4)

        except Exception as e:
            self.logger.error(f"Index creation error: {str(e)}", exc_info=True)